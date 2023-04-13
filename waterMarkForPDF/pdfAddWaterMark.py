import os
import io
import base64
import asyncio
from pyppeteer import launch
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def add_watermark(input_ppt, watermark_image, output_ppt):
    # 加载PPT文件
    prs = Presentation(input_ppt)

    # 加载水印图片
    watermark_img = Image.open(watermark_image)

    # 设置水印尺寸
    watermark_width = Inches(7)
    watermark_height = Inches(7)

    # 计算根据幻灯片大小和水印大小所需的水印行数和列数
    slide_width = prs.slide_width.emu
    slide_height = prs.slide_height.emu
    columns = int(slide_width / watermark_width.emu) + 1
    rows = int(slide_height / watermark_height.emu) + 1

    # 在每个幻灯片上平铺水印
    for slide in prs.slides:
        for r in range(rows):
            for c in range(columns):
                x_pos = c * watermark_width.emu
                y_pos = r * watermark_height.emu
                slide.shapes.add_picture(watermark_image, x_pos, y_pos, width=watermark_width, height=watermark_height)

    # 保存带水印的PPT文件
    prs.save(output_ppt)

def save_ppt_as_images(input_ppt):
    prs = Presentation(input_ppt)
    image_data_list = []

    for i, slide in enumerate(prs.slides):
        # Create a file-like buffer to receive the image data
        image_data = io.BytesIO()

        # Save the slide as an image
        image = Image.frombytes("RGB", (prs.slide_width, prs.slide_height), slide.shapes[0].image.blob)
        image.save(image_data, "JPEG")

        # Convert the image data to base64
        image_base64 = base64.b64encode(image_data.getvalue()).decode("utf-8")

        image_data_list.append(f'<img src="data:image/jpeg;base64,{image_base64}" style="width: 100%; page-break-after: always;">')

    return "\n".join(image_data_list)

async def convert_to_pdf(ppt_file, pdf_file):
    browser = await launch()
    page = await browser.newPage()

    await page.setViewport({"width": 1280, "height": 720})

    images_html = save_ppt_as_images(ppt_file)
    html_content = f"""
    <!doctype html>
    <html>
        <head>
            <style>
                body {{
                    margin: 0;
                    padding: 0;
                }}
            </style>
        </head>
        <body>
            {images_html}
        </body>
    </html>
    """
    await page.setContent(html_content, {"waitUntil": "networkidle0"})
    print(">>>>")

    await page.pdf({"path": pdf_file, "format": "A4", "printBackground": True})

    await browser.close()

if __name__ == "__main__":
    base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets')
    input_ppt = os.path.join(base_path, 'input.pptx')
    watermark_image = os.path.join(base_path, 'watermark.png')
    output_ppt = os.path.join(base_path, 'output.pptx')
    output_pdf = os.path.join(base_path, 'output.pdf')

    # 1. 在 PPT 上添加水印E
    add_watermark(input_ppt, watermark_image, output_ppt)

    # 2. 将带有水印的 PPT 转换为 PDF
    asyncio.get_event_loop().run_until_complete(convert_to_pdf(output_ppt, output_pdf))

    print("完成。已将带有水印的PPT转换为PDF。")