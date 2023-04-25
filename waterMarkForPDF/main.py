import os
import io
import base64
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import aspose.slides as slides

def add_watermark(input_ppt, watermark_image, output_ppt):
    prs = Presentation(input_ppt)
    watermark_img = Image.open(watermark_image)
    watermark_width = Inches(1)
    watermark_height = Inches(1)
    slide_width = prs.slide_width.emu
    slide_height = prs.slide_height.emu
    columns = int(slide_width / watermark_width.emu) + 1
    rows = int(slide_height / watermark_height.emu) + 1

    for slide in prs.slides:
        for r in range(rows):
            for c in range(columns):
                x_pos = c * watermark_width.emu
                y_pos = r * watermark_height.emu
                slide.shapes.add_picture(watermark_image, x_pos, y_pos, width=watermark_width, height=watermark_height)

    prs.save(output_ppt)

def convert_to_pdf_aspose(pptx_file, pdf_file):
    with slides.Presentation(pptx_file) as presentation:
        presentation.save(pdf_file, slides.export.SaveFormat.PDF)

if __name__ == "__main__":
    base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets')
    input_folder = os.path.join(base_path, 'input')
    watermark_image = os.path.join(base_path, 'watermark.png')
    output_folder = os.path.join(base_path, 'output')

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for ppt_file in os.listdir(input_folder):
        if ppt_file.endswith('.pptx'):
            input_ppt = os.path.join(input_folder, ppt_file)
            output_ppt = os.path.join(output_folder, ppt_file)
            output_pdf = os.path.join(output_folder, ppt_file.replace('.pptx', '.pdf'))
            
            # 1. 在 PPT 上添加水印
            add_watermark(input_ppt, watermark_image, output_ppt)

            # 2. 将带有水印的 PPT 转换为 PDF
            # convert_to_pdf_aspose(output_ppt, output_pdf)

            print(f"完成。已将带有水印的 PPT 转换为 PDF，并保存到：{output_pdf}")
